"""Helpers de formatação para os slides de projeção (python-pptx):
16:9, texto grande e legível à distância, sem elementos decorativos —
um slide por seção do mesmo conteúdo usado na Pauta/Ata em DOCX.
"""

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

PRETO = RGBColor(0x00, 0x00, 0x00)
CINZA = RGBColor(0x40, 0x40, 0x40)

FONTE = 'Calibri'
TAM_TITULO_CAPA = Pt(44)
TAM_SUBTITULO_CAPA = Pt(24)
TAM_TITULO_SECAO = Pt(36)
TAM_CORPO = Pt(28)

LARGURA = Inches(13.333)
ALTURA = Inches(7.5)


def nova_apresentacao():
    prs = Presentation()
    prs.slide_width = LARGURA
    prs.slide_height = ALTURA
    return prs


def _slide_em_branco(prs):
    layout_branco = prs.slide_layouts[6]
    return prs.slides.add_slide(layout_branco)


def _formatar_paragrafo(paragrafo, texto, tamanho, negrito=False, cor=PRETO, alinhamento=None):
    run = paragrafo.add_run()
    run.text = texto
    run.font.size = tamanho
    run.font.bold = negrito
    run.font.name = FONTE
    run.font.color.rgb = cor
    if alinhamento is not None:
        paragrafo.alignment = alinhamento
    return run


def slide_capa(prs, titulo, subtitulo):
    slide = _slide_em_branco(prs)

    caixa_titulo = slide.shapes.add_textbox(Inches(0.8), Inches(2.7), Inches(11.73), Inches(1.3))
    tf = caixa_titulo.text_frame
    tf.word_wrap = True
    _formatar_paragrafo(tf.paragraphs[0], titulo, TAM_TITULO_CAPA, negrito=True, alinhamento=PP_ALIGN.CENTER)

    caixa_sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.73), Inches(0.8))
    tf2 = caixa_sub.text_frame
    tf2.word_wrap = True
    _formatar_paragrafo(tf2.paragraphs[0], subtitulo, TAM_SUBTITULO_CAPA, cor=CINZA, alinhamento=PP_ALIGN.CENTER)

    return slide


def _slide_com_titulo(prs, titulo):
    slide = _slide_em_branco(prs)
    caixa_titulo = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    tf = caixa_titulo.text_frame
    tf.word_wrap = True
    _formatar_paragrafo(tf.paragraphs[0], titulo, TAM_TITULO_SECAO, negrito=True)
    return slide


def slide_conteudo(prs, titulo, conteudo, checklist=False):
    """conteudo: string (um parágrafo) ou lista (um bullet por item)."""
    slide = _slide_com_titulo(prs, titulo)
    caixa_corpo = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.3))
    tf = caixa_corpo.text_frame
    tf.word_wrap = True

    itens = conteudo if isinstance(conteudo, (list, tuple)) else [conteudo]
    for i, item in enumerate(itens):
        paragrafo = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if checklist:
            prefixo = '☐ '
        elif len(itens) > 1:
            prefixo = '• '
        else:
            prefixo = ''
        _formatar_paragrafo(paragrafo, f'{prefixo}{item}', TAM_CORPO)
        paragrafo.space_after = Pt(16)

    return slide


def slide_participantes(prs, participantes):
    slide = _slide_com_titulo(prs, 'Participantes')
    caixa_corpo = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.3))
    tf = caixa_corpo.text_frame
    tf.word_wrap = True

    for i, participante in enumerate(participantes):
        paragrafo = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(participante, dict):
            marca = '✓' if participante.get('presente', True) else '✗'
            texto = f"{marca} {participante['nome']}"
        else:
            texto = f'• {participante}'
        _formatar_paragrafo(paragrafo, texto, TAM_CORPO)
        paragrafo.space_after = Pt(10)

    return slide
